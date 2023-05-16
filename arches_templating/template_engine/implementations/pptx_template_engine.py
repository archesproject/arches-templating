from base64 import b64decode
import copy
from io import BytesIO
import re
from typing import List, Tuple

from arches_templating.template_engine.template_engine_factory import TemplateEngineFactory
from arches_templating.template_engine.template_tag_type import TemplateTagType
from arches_templating.template_engine.template_engine import TemplateEngine
from arches_templating.template_engine.template_tag import TemplateTag
from pptx import Presentation
from pptx.shapes.graphfrm import GraphicFrame
from pptx.table import _Cell, _Row

@TemplateEngineFactory.register('pptx')
class PptxTemplateEngine(TemplateEngine):

    def extract_regex_matches(self, template) -> List[Tuple]:
        self.presentation = Presentation(template)
        parsed_tags: List[Tuple] = []
        parsed_tags += self.iterate_over_container(self.presentation.slides)
        return parsed_tags
          # should match <arches: node_alias>

    def match_text(self, container, parent, cell=None):
        parsed_tags: List[Tuple] = []
        for match in re.findall(self.regex, container.text):
            parsed_tags.append((match, {"container": container, "parent": parent, "cell": None}))
        return parsed_tags

    def iterate_over_container(self, container, parent=None):
        parsed_tags: List[Tuple] = []
        if hasattr(container, '__iter__'):
            for s in container:
                for shape in s.shapes:
                    if type(shape) is not GraphicFrame:
                        parsed_tags += self.match_text(shape, container)
                        
                    else:
                        try:
                            table = shape.table
                            row_length = len(table.rows)
                            column_length = len(table.columns)
                            current_row = 0
                            while current_row < row_length:
                                current_column = 0
                                while current_column < column_length:
                                    current_cell = table.cell(current_row, current_column)
                                    parsed_tags += self.iterate_over_container(current_cell, table)
                                    current_column += 1
                                current_row += 1
                            pass
                        except AttributeError:
                            pass # ok to pass; happens if the shape doesn't have a table - in this case, don't render subtags - we don't know how.
        elif type(container) is _Cell:
            parsed_tags += self.match_text(container, parent)
        return parsed_tags

    def delete_element(element):
        elem = element._element
        elem.getparent().remove(elem)
        elem._element = None

    def remove_row(table, row_index):
        """
        Removes specified *row* (e.g. ``table.rows.remove(table.rows[0])``).
        """
        row = table.rows[row_index]
        table._tbl.remove(row._tr)

    def add_row(table, index=-1):
        """
        Duplicates last row to keep formatting and resets it's cells text_frames
        (e.g. ``row = table.rows.add_row()``).
        Returns new |_Row| instance.
        """
        new_row = copy.deepcopy(table._tbl.tr_lst[2])  # copies last row element

        for tc in new_row.tc_lst:
            cell = _Cell(tc, new_row.tc_lst)
            cell.text = ''

        if index == -1:
            index = len(table._tbl) % index

        table._tbl.insert(index, new_row)

        return _Row(new_row, table)

    def replace_tags(self, tags:List[TemplateTag]):
        for tag in tags:
            block = tag.optional_keys['container']
            if tag.type == TemplateTagType.CONTEXT:                    
                block.text = block.text.replace(tag.raw, "")
                end_block = tag.end_tag.optional_keys['container']
                end_block.text = end_block.text.replace(tag.end_tag.raw, "")

                if tag.has_rows:
                    column = 0
                    # this is ugly, but way more efficient than the alternative
                    parent_table = tag.context_children_template[-1].optional_keys["parent"]
                    
                    current_row = PptxTemplateEngine.add_row(parent_table)


                    for child in tag.children:
                        if child.type == TemplateTagType.ROWEND:
                            column = -1
                            current_row = PptxTemplateEngine.add_row(parent_table)
                        elif child.type == TemplateTagType.VALUE:
                            # grab any borders from the original cell copy them to the new cell.
                            #template_block = tag.context_children_template[column].optional_keys["container"]

                            paragraph = current_row.cells[column].text_frame.add_paragraph()
                            paragraph.text = "" if child.value == None else child.value
                        column += 1
                    
                    lead_matches_len = len(self.regex.findall(parent_table.cell(0,0).text))

                    PptxTemplateEngine.remove_row(parent_table, 2)
                    if lead_matches_len > 1:
                        primary_cell = parent_table.cell(0,0)
                        primary_cell.text = primary_cell.text.replace(tag.raw, primary_cell.text)
                    else:
                        PptxTemplateEngine.remove_row(parent_table, 0)
                    
                else:
                    self.replace_tags(tag.children)

            elif tag.type == TemplateTagType.VALUE:
                block.text = block.text.replace(tag.raw, tag.value)
            elif tag.type == TemplateTagType.IMAGE:
                block._parent.add_picture(BytesIO(b64decode(re.sub("data:image/jpeg;base64,", '', tag.value))), block.top, block.left, block.width, block.height)
                PptxTemplateEngine.delete_element(block)
            elif tag.type == TemplateTagType.IF:
                if tag.render:
                    block.text = block.text.replace(tag.raw, "")
                    end_block = tag.end_tag.optional_keys['container']
                    end_block.text = end_block.text.replace(tag.end_tag.raw, "")
                    self.replace_tags(tag.children)
                else:
                    all_elements = []
                    found_if_start = False
                    found_if_end = False
                    for item in self.iter_block_items(block._parent):
                        if item._element == block._element:
                            found_if_start = True
                        if item._element == tag.end_tag.optional_keys["docxBlock"]._element:
                            found_if_end = True
                            try:
                                item.text = ""
                            except AttributeError: # ok to skip if text is not found
                                pass
                        if found_if_start and not found_if_end:
                            try:
                                item.text = ""
                            except AttributeError: # ok to skip if text is not found
                                pass

    def create_file(self, tags:List[TemplateTag], template):
        bytestream = BytesIO()
        mime = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        self.replace_tags(tags)
        self.presentation.save(bytestream)
        return (bytestream, mime, "test.pptx")